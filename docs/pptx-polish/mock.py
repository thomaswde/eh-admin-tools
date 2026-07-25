"""System Health PPTX polish mock-up.

Draws 4 representative slides in Light and Dark themes using only primitives
that PptxGenJS can also emit (rects, rounded rects, lines, ellipses, text,
images). Nothing here relies on python-pptx-only capability.
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw

W, H = 13.333, 7.5
M = 0.7
FONT = "Arial"

LIGHT = dict(
    name="LIGHT",
    bg="FFFFFF", ink="261F63", body="3D3760", muted="7B7799",
    rule="E5E4ED", track="EFEEF4", band="F8F8FB",
    crit="EC0889", warn="F05918", ok="00AAEF", gray="C3C1D0",
    logo="assets/eh-logo-color.png",
)
DARK = dict(
    name="DARK",
    bg="141221", ink="FFFFFF", body="D5D2E6", muted="8F8BAD",
    rule="2E2A45", track="242038", band="1C1930",
    crit="FF3BA5", warn="FF8043", ok="2FBDF5", gray="4E4968",
    logo="assets/eh-logo-white.png",
)

SAPPHIRE = (0x26, 0x1F, 0x63)
PLUM = (0x7F, 0x28, 0x54)
LIME = "DAED43"


# ---------------------------------------------------------------- helpers
def rgb(h):
    return RGBColor.from_string(h)


def gradient_png(path, w=2000, h=1125):
    """Sapphire -> Plum diagonal, with restrained concentric ring motif."""
    img = Image.new("RGB", (w, h))
    px = img.load()
    for y in range(h):
        for x in range(0, w, 4):
            t = (x / w) * 0.72 + (y / h) * 0.28
            c = tuple(int(SAPPHIRE[i] + (PLUM[i] - SAPPHIRE[i]) * t) for i in range(3))
            for dx in range(4):
                if x + dx < w:
                    px[x + dx, y] = c
    ring = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(ring)
    cx, cy = int(w * 0.80), int(h * 0.42)
    for i in range(16):
        r = int(w * 0.055) + i * int(w * 0.028)
        a = max(0, 26 - i)
        d.ellipse([cx - r, cy - int(r * 0.78), cx + r, cy + int(r * 0.78)],
                  outline=(255, 255, 255, a), width=2)
    img = Image.alpha_composite(img.convert("RGBA"), ring).convert("RGB")
    img.save(path, "PNG")
    return path


def text(slide, x, y, w, h, s, size=12, color="261F63", bold=False,
         align=PP_ALIGN.LEFT, space=None, anchor=MSO_ANCHOR.TOP, line=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(s).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        if space:
            # character spacing (tracking)
            r.font._rPr.set("spc", str(int(space * 100)))
    return tb


def rect(slide, x, y, w, h, fill=None, line_col=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE, adj=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if adj is not None:
        try:
            sh.adjustments[0] = adj
        except Exception:
            pass
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
    else:
        sh.fill.background()
    if line_col:
        sh.line.color.rgb = rgb(line_col)
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    from pptx.oxml.ns import qn
    # drop the theme style reference; its effectRef re-applies a drop shadow
    for st in sh._element.findall(qn('p:style')):
        sh._element.remove(st)
    spPr = sh._element.spPr
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    spPr.append(spPr.makeelement(qn('a:effectLst'), {}))
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return sh


def hline(slide, x, y, w, color, width=0.75):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(width)
    return ln


def logo(slide, t, x, y, w):
    slide.shapes.add_picture(t["logo"], Inches(x), Inches(y), width=Inches(w))


# ---------------------------------------------------------------- chrome
def blank(prs, t):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = rect(s, -0.02, -0.02, W + 0.04, H + 0.04, fill=t["bg"])
    return s


def body_slide(prs, t, title, sub, page, kicker=None, kicker_col=None):
    s = blank(prs, t)
    y = 0.52
    if kicker:
        text(s, M, y, 6.0, 0.2, kicker.upper(), 9,
             kicker_col or t["crit"], True, space=1.6)
        y += 0.30
    text(s, M, y, 11.4, 0.5, title, 26, t["ink"], True)
    if sub:
        text(s, M, y + 0.50, 11.4, 0.26, sub, 11.5, t["muted"])
    hline(s, M, 7.02, W - 2 * M, t["rule"], 0.75)
    text(s, M, 7.11, 6.0, 0.18, "ExtraHop  ·  System Health Review", 8, t["muted"])
    text(s, 10.2, 7.11, 1.05, 0.18, str(page), 8, t["muted"], align=PP_ALIGN.RIGHT)
    logo(s, t, 11.63, 7.06, 1.0)
    return s


# ---------------------------------------------------------------- slide 1
def cover(prs, t, grad):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(grad, Inches(-0.02), Inches(-0.02),
                         width=Inches(W + 0.04), height=Inches(H + 0.04))
    s.shapes.add_picture("assets/eh-logo-white.png", Inches(0.86), Inches(0.66),
                         width=Inches(1.95))

    # Lozenge enclosure, bleeding off the left edge (brand shape signature)
    rect(s, -1.30, 2.44, 9.15, 2.06, fill=None, line_col="FFFFFF", line_w=0.9,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)

    text(s, 0.86, 2.78, 7.4, 0.62, "System Health Review", 38, "FFFFFF", True)
    text(s, 0.88, 3.56, 7.0, 0.34, "ExtraHop SE Organization", 17, LIME, True)

    hline(s, 0.88, 6.10, 4.2, "FFFFFF", 0.75)
    meta = [("REPORT WINDOW", "Last 7 days"),
            ("PREPARED BY", "Thomas Smith"),
            ("GENERATED", "25 Jul 2026")]
    for i, (k, v) in enumerate(meta):
        x = 0.88 + i * 2.55
        text(s, x, 6.28, 2.4, 0.18, k, 8, "FFFFFF", True, space=1.4)
        text(s, x, 6.52, 2.4, 0.22, v, 12, "FFFFFF")
    return s


# ---------------------------------------------------------------- slide 2
def glance(prs, t):
    s = body_slide(prs, t, "Fleet health at a glance",
                   "78 sensors · Last 7 days · Peak 1-hour averages", 2)

    # ---- verdict
    text(s, M, 1.52, 11.9, 0.32,
         "62 of 78 sensors returned no data. Capacity is not the constraint — reachability is.",
         15.5, t["ink"], True)

    # ---- stacked fleet status bar
    segs = [("Reporting", 14, t["ok"]), ("Needs attention", 2, t["crit"]),
            ("No data returned", 62, t["gray"])]
    total = sum(c for _, c, _ in segs)
    bar_y, bar_h, bar_w = 2.12, 0.34, W - 2 * M
    x = M
    for label, count, col in segs:
        w = max(bar_w * count / total, 0.22)
        rect(s, x, bar_y, w - 0.035, bar_h, fill=col)
        x += w
    # legend — evenly spaced so short segments stay readable
    for i, (label, count, col) in enumerate(segs):
        x = M + i * 2.75
        rect(s, x, 2.615, 0.1, 0.1, fill=col)
        text(s, x + 0.20, 2.565, 2.4, 0.2, f"{count}  {label}", 10.5, t["body"])

    hline(s, M, 3.22, W - 2 * M, t["rule"])

    # ---- stat row
    stats = [("78", "Sensors in fleet", "returned by the appliance API", t["ink"]),
             ("14", "Reporting data", "18% of the fleet", t["ink"]),
             ("2", "At or over capacity", "advanced analysis tier", t["crit"]),
             ("1.4M", "Trigger drops", "total across the window", t["crit"])]
    for i, (val, label, note, col) in enumerate(stats):
        x = M + i * 3.06
        text(s, x, 3.52, 2.8, 0.62, val, 38, col, True)
        text(s, x, 4.22, 2.8, 0.22, label, 11.5, t["ink"], True)
        text(s, x, 4.48, 2.8, 0.2, note, 9.5, t["muted"])
        hline(s, x, 4.82, 0.62, col, 2.0)

    hline(s, M, 5.28, W - 2 * M, t["rule"])

    # ---- composition chips
    text(s, M, 5.52, 3.0, 0.2, "FLEET COMPOSITION", 9, t["muted"], True, space=1.5)
    chips = [("EDA1100V_TRACE", 48), ("IDS1280V", 15), ("EDA1100V", 9),
             ("EFC1290V", 3), ("EFC1291V", 3)]
    x = M
    for name, n in chips:
        w = 0.115 * len(name) + 0.62
        rect(s, x, 5.82, w, 0.34, fill=t["band"], line_col=t["rule"],
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
        text(s, x + 0.22, 5.90, w - 0.44, 0.2, name, 10, t["body"])
        text(s, x + w - 0.55, 5.90, 0.34, 0.2, f"×{n}", 10, t["ink"], True,
             align=PP_ALIGN.RIGHT)
        x += w + 0.14

    text(s, M, 6.44, 11.9, 0.22,
         "Counts reflect sensors that returned data. Unavailable statistics are held as "
         "missing rather than zero; per-sensor collection status is in the appendix.",
         10, t["muted"])
    return s


# ---------------------------------------------------------------- slide 3
def attention(prs, t):
    s = body_slide(prs, t, "Sensors that need attention",
                   "Ranked by severity · grouped by condition", 3,
                   kicker="4 sensors · 2 critical")

    cols = [("SENSOR", M + 0.20, 2.55), ("MODEL", 3.62, 1.75),
            ("CONDITION", 5.48, 2.35), ("EVIDENCE", 7.95, 4.65)]
    hy = 1.72
    for label, x, w in cols:
        text(s, x, hy, w, 0.2, label, 8.5, t["muted"], True, space=1.4)
    hline(s, M, hy + 0.24, W - 2 * M, t["rule"])

    groups = [
        ("crit", "CRITICAL", [
            ("robertb-gcp-lab1", "EDA1100V_TRACE", "Advanced analysis full",
             "250 / 250 devices · 6 more waiting in Discovery · 361,605 trigger drops"),
            ("robertb-sensor-3", "EDA1100V_TRACE", "Trigger drops",
             "800,133 drops · trigger load peaked at 45% of available cycles"),
        ]),
        ("warn", "WARNING", [
            ("eda.Pai.CyberRange", "EDA1100V", "Trigger drops",
             "265,282 drops · trigger load peaked at 4% of available cycles"),
            ("SE-O365", "EDA1100V", "Stale synchronization",
             "Last synchronized 15 Jan 2026 · 191 days ago"),
        ]),
    ]

    y = hy + 0.36
    rh = 0.58
    band = 0
    for sev, label, rows in groups:
        text(s, M + 0.20, y + 0.02, 2.0, 0.2, label, 8, t[sev], True, space=1.4)
        y += 0.26
        for name, model, cond, detail in rows:
            if band % 2 == 0:
                rect(s, M, y, W - 2 * M, rh, fill=t["band"])
            rect(s, M, y, 0.055, rh, fill=t[sev])
            text(s, M + 0.20, y, 2.55, rh, name, 11.5, t["ink"], True,
                 anchor=MSO_ANCHOR.MIDDLE)
            text(s, 3.62, y, 1.75, rh, model, 10.5, t["muted"],
                 anchor=MSO_ANCHOR.MIDDLE)
            text(s, 5.48, y, 2.35, rh, cond, 11, t[sev], True,
                 anchor=MSO_ANCHOR.MIDDLE)
            text(s, 7.95, y, 4.70, rh, detail, 10.5, t["body"],
                 anchor=MSO_ANCHOR.MIDDLE)
            y += rh
            band += 1
        y += 0.16

    # roll-up band — the 62 no-data sensors collapse to one line
    y += 0.10
    rect(s, M, y, W - 2 * M, 0.84, fill=t["band"], line_col=t["rule"],
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.16)
    rect(s, M, y, 0.055, 0.84, fill=t["gray"])
    text(s, M + 0.34, y + 0.15, 8.6, 0.26, "62 sensors returned no data", 13,
         t["ink"], True)
    text(s, M + 0.34, y + 0.45, 8.8, 0.24,
         "61 unreachable · 1 requires additional configuration. "
         "No utilization conclusions are drawn for these sensors.",
         10.5, t["muted"])
    text(s, 10.0, y + 0.29, 2.63, 0.24, "Full list in appendix", 10.5,
         t["muted"], align=PP_ALIGN.RIGHT)
    return s


# ---------------------------------------------------------------- slide 4
def capacity(prs, t):
    s = body_slide(prs, t, "Packet rate headroom",
                   "Peak 1-hour average against model capacity · EDA1100V_TRACE · rated 140K p/s",
                   4, kicker="Capacity · 1 of 3", kicker_col=t["muted"])

    def clip(name, n=30):
        return name if len(name) <= n else name[:n - 1] + "…"

    x0, x1 = 3.35, 10.35          # plot area: x1 == 100% of model capacity
    span = x1 - x0
    top = 1.90
    rh = 0.44

    bars = [("robertb-gcp-lab1", 0.21, "29K p/s"),
            ("robertb-sensor-3", 0.045, "6.3K p/s"),
            ("7_22_26_SalkeId_Ultra", 0.005, "712 p/s"),
            ("eda.Casillas3.us-east-2", 0.0017, "238 p/s"),
            ("eda.Pankaj-New-CR.ap-southeast-2", 0.0015, "215 p/s"),
            ("ariel-aws-ultra-eda", 0.0009, "120 p/s"),
            ("eda.natejohnson.us-east-2", 0.0008, "116 p/s"),
            ("ArielSmi", 0.0008, "107 p/s"),
            ("AWS-SE360", 0.0005, "63 p/s"),
            ("eda.jasonj-360-range-2026.us-east-2", 0.0001, "12 p/s")]

    bottom = top + rh * len(bars)

    # 80% guide only — the track already ends at 100% of capacity
    gx = x0 + span * 0.80
    ln = s.shapes.add_connector(1, Inches(gx), Inches(top - 0.06),
                                Inches(gx), Inches(bottom + 0.02))
    ln.line.color.rgb = rgb(t["warn"])
    ln.line.width = Pt(0.75)
    ln.line.dash_style = 4
    text(s, gx - 0.40, top - 0.28, 0.8, 0.18, "80%", 7.5, t["warn"],
         True, space=1.0, align=PP_ALIGN.CENTER)
    text(s, x1 - 1.8, top - 0.28, 1.8, 0.18, "MODEL CAPACITY", 7.5, t["muted"],
         True, space=1.0, align=PP_ALIGN.RIGHT)

    y = top
    for name, frac, val in bars:
        sev = "crit" if frac >= 1 else "warn" if frac >= 0.8 else "ok"
        text(s, M, y + 0.055, 2.45, 0.26, clip(name), 10, t["body"],
             align=PP_ALIGN.RIGHT)
        rect(s, x0, y + 0.065, span, 0.24, fill=t["track"])
        rect(s, x0, y + 0.065, max(span * frac, 0.035), 0.24, fill=t[sev])
        pct = f"{frac*100:.0f}%" if frac >= 0.01 else "<1%"
        text(s, x1 + 0.20, y + 0.055, 0.66, 0.26, pct, 10, t["ink"], True,
             align=PP_ALIGN.RIGHT)
        text(s, x1 + 0.96, y + 0.055, 1.15, 0.26, val, 10, t["muted"],
             align=PP_ALIGN.RIGHT)
        y += rh

    hline(s, M, y + 0.22, W - 2 * M, t["rule"])
    text(s, M, y + 0.36, 8.8, 0.22,
         "10 of 48 EDA1100V_TRACE sensors reported packet data · "
         "38 returned no data and are excluded rather than shown as zero.",
         10, t["muted"])
    text(s, 9.2, y + 0.36, 3.45, 0.22, "Sorted by percent of model capacity",
         10, t["muted"], align=PP_ALIGN.RIGHT)
    return s


# ---------------------------------------------------------------- build
def build(theme, out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    cover(prs, theme, "grad.png")
    glance(prs, theme)
    attention(prs, theme)
    capacity(prs, theme)
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    gradient_png("grad.png")
    build(LIGHT, "mockup-light.pptx")
    build(DARK, "mockup-dark.pptx")
